# -*- coding: utf-8 -*-
"""
Генерация презентации «Полугодовая премия ИТ-работников (2 полугодие 2026)».

КАК ЗАПУСТИТЬ (Windows, командная строка / PowerShell):
    1) Установить библиотеку (один раз):
           pip install python-pptx
    2) Запустить скрипт:
           python generate_presentation.py
       либо указать свою папку для сохранения:
           python generate_presentation.py "C:\\Users\\Имя\\Documents\\Презентации"

Если папку не указать — файл сохранится рядом со скриптом.
"""

import os
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    sys.exit(
        "Не установлена библиотека python-pptx.\n"
        "Установите её командой:  pip install python-pptx\n"
        "и запустите скрипт снова."
    )

# ----------------------------------------------------------------------------
# НАСТРОЙКИ: куда сохранять и как назвать файл
# ----------------------------------------------------------------------------
# Папка по умолчанию (если не передать аргумент в командной строке).
# Можно вписать сюда свой путь, например: r"C:\Users\Имя\Documents\Презентации"
DEFAULT_OUTPUT_DIR = "."
OUTPUT_FILENAME = "Полугодовая_премия_ИТ-работников_2ПГ_2026.pptx"

# ----------------------------------------------------------------------------
# Палитра (корпоративный, «финансово-доверительный» стиль)
# ----------------------------------------------------------------------------
DARK   = RGBColor(0x0E, 0x3A, 0x36)   # глубокий тёмно-зелёный (титул, акцентный слайд)
TEAL   = RGBColor(0x1B, 0x99, 0x8B)   # бирюзово-зелёный
TEAL_D = RGBColor(0x12, 0x6B, 0x60)   # тёмная бирюза (контраст на светлом)
GOLD   = RGBColor(0xE2, 0xA4, 0x3B)   # золотой акцент (деньги, ключевые цифры)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
INK    = RGBColor(0x1B, 0x26, 0x24)   # почти чёрный текст
MUTED  = RGBColor(0x5C, 0x6E, 0x6B)   # приглушённый серо-зелёный
TINT   = RGBColor(0xEA, 0xF3, 0xF1)   # светлая бирюзовая заливка карточек
TINT2  = RGBColor(0xF4, 0xF9, 0xF7)   # очень светлая заливка
GOLDT  = RGBColor(0xFA, 0xF0, 0xDC)   # светлый золотой фон
LINE   = RGBColor(0xD3, 0xE2, 0xDE)   # рамка карточек

EMU = 914400  # 1 дюйм


# ----------------------------------------------------------------------------
# Вспомогательные функции построения слайдов
# ----------------------------------------------------------------------------
def blank_slide(prs):
    """Пустой слайд (без плейсхолдеров)."""
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_bg(slide, color):
    """Сплошной фон слайда."""
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _no_shadow(shape):
    # отключаем тень по умолчанию
    shape.shadow.inherit = False


def rect(slide, l, t, w, h, fill, line=None, line_w=1.0, rounded=False, radius=0.06):
    """Прямоугольник / скруглённый прямоугольник."""
    shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shp_type, Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    _no_shadow(shp)
    if rounded:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    return shp


def oval(slide, l, t, d, fill, h=None):
    """Круг (или эллипс, если задана высота h)."""
    shp = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(d), Inches(h if h else d)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    _no_shadow(shp)
    return shp


def _set_hang(p, marL=0.26):
    """Висячий отступ: перенесённые строки выравниваются под текстом, не под маркером."""
    pPr = p._p.get_or_add_pPr()
    pPr.set("marL", str(int(marL * EMU)))
    pPr.set("indent", str(int(-marL * EMU)))


def text(slide, s, l, t, w, h, size=16, color=INK, bold=False, italic=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri",
         line_spacing=1.0, space_after=0):
    """Текстовый блок. Многострочность через '\\n'."""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, line in enumerate(s.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        if space_after:
            p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = line
        f = run.font
        f.size = Pt(size)
        f.bold = bold
        f.italic = italic
        f.name = font
        f.color.rgb = color
    return tb


def bullets(slide, items, l, t, w, h, size=15, color=INK, glyph="•",
            glyph_color=None, gap=8, line_spacing=1.06, font="Calibri"):
    """Маркированный список с висячим отступом."""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(gap)
        _set_hang(p)
        rb = p.add_run()
        rb.text = glyph + "  "
        rb.font.size = Pt(size)
        rb.font.name = font
        rb.font.bold = True
        rb.font.color.rgb = glyph_color or GOLD
        rt = p.add_run()
        rt.text = it
        rt.font.size = Pt(size)
        rt.font.name = font
        rt.font.color.rgb = color
    return tb


def badge(slide, l, t, d, fill, glyph, glyph_color=WHITE, glyph_size=18,
          font="Calibri", bold=True):
    """Круглый значок с символом/цифрой (фирменный мотив)."""
    c = oval(slide, l, t, d, fill)
    tf = c.text_frame
    tf.word_wrap = False
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = glyph
    r.font.size = Pt(glyph_size)
    r.font.bold = bold
    r.font.name = font
    r.font.color.rgb = glyph_color
    return c


def slide_title(slide, title, kicker=None, color=DARK):
    """Заголовок контентного слайда + фирменная точка-мотив слева."""
    oval(slide, 0.62, 0.62, 0.2, GOLD)
    text(slide, title, 0.98, 0.46, 11.6, 0.7, size=29, bold=True,
         color=color, font="Cambria", anchor=MSO_ANCHOR.MIDDLE)
    if kicker:
        text(slide, kicker, 0.98, 1.12, 11.6, 0.4, size=13, color=MUTED, font="Calibri")


# ----------------------------------------------------------------------------
# Слайды
# ----------------------------------------------------------------------------
def s01_title(prs):
    s = blank_slide(prs)
    set_bg(s, DARK)
    # геометрический мотив — перекрывающиеся круги в правом нижнем углу
    oval(s, 9.7, 3.9, 5.2, TEAL)
    oval(s, 11.3, 5.4, 2.7, GOLD)
    oval(s, 8.9, 5.9, 1.3, WHITE)

    text(s, "Июнь 2026", 0.95, 1.5, 6, 0.4, size=14, color=GOLD,
         bold=True, font="Calibri")
    text(s, "Полугодовая премия\nИТ-работников", 0.95, 2.1, 9.0, 2.0,
         size=46, bold=True, color=WHITE, font="Cambria", line_spacing=1.02)
    text(s, "2 полугодие 2026", 0.95, 4.15, 8, 0.6, size=24, color=GOLD,
         bold=True, font="Calibri")
    text(s, "БЕ Технологии   ·   Цифровые Платформы   ·   БЕ Роботикс",
         0.97, 5.05, 8.5, 0.4, size=14, color=TINT, font="Calibri")
    text(s, "ИТ-работники · бэнды M, LM, S", 0.97, 5.5, 8.5, 0.4,
         size=14, color=MUTED, font="Calibri")
    return s


def s02_contents(prs):
    s = blank_slide(prs)
    set_bg(s, WHITE)
    slide_title(s, "Содержание", "Навигация по разделам")

    items = [
        ("01", "Как работает система", "Две шкалы, группы, принцип"),
        ("02", "Как формируется оценка", "От КПЭ до калибровки"),
        ("03", "Коэффициент фонда (Кф)", "Корректировки на 2 полугодие"),
        ("04", "Практический пример", "Пошаговый расчёт премии"),
        ("05", "Гарантии для вас", "Четыре механизма защиты"),
        ("06", "Сроки и условия выплаты", "Когда и при каких условиях"),
        ("07", "Частые вопросы", "Декрет и калибровка"),
        ("08", "Ваши следующие шаги", "Чек-лист и контакты"),
    ]
    # две колонки по 4 пункта
    col_x = [0.95, 7.0]
    top0 = 1.95
    row_h = 1.18
    for i, (num, title, sub) in enumerate(items):
        col = 0 if i < 4 else 1
        row = i % 4
        x = col_x[col]
        y = top0 + row * row_h
        badge(s, x, y, 0.62, TEAL, num, WHITE, glyph_size=16, font="Cambria")
        text(s, title, x + 0.85, y - 0.02, 4.5, 0.45, size=16, bold=True,
             color=INK, font="Calibri", anchor=MSO_ANCHOR.MIDDLE)
        text(s, sub, x + 0.85, y + 0.36, 4.5, 0.4, size=12, color=MUTED,
             font="Calibri")
    return s


def s03_how_it_works(prs):
    s = blank_slide(prs)
    set_bg(s, WHITE)
    slide_title(s, "Как работает система", "Две шкалы · группы · единые правила")

    # Карточка: Шкала №1
    rect(s, 0.95, 1.95, 5.55, 2.05, TINT, rounded=True, radius=0.05)
    badge(s, 1.3, 2.3, 0.6, TEAL_D, "1", WHITE, glyph_size=20, font="Cambria")
    text(s, "Шкала №1", 2.1, 2.32, 4.0, 0.5, size=18, bold=True, color=INK,
         font="Calibri")
    text(s, "Базовая шкала оценки эффективности. Переводит итоговую оценку "
            "1–5 в процент выполнения премии.",
         1.3, 3.05, 4.9, 0.9, size=13.5, color=INK, font="Calibri",
         line_spacing=1.08)

    # Карточка: Шкала №2
    rect(s, 6.8, 1.95, 5.55, 2.05, TINT, rounded=True, radius=0.05)
    badge(s, 7.15, 2.3, 0.6, GOLD, "2", WHITE, glyph_size=20, font="Cambria")
    text(s, "Шкала №2", 7.95, 2.32, 4.0, 0.5, size=18, bold=True, color=INK,
         font="Calibri")
    text(s, "Расширенная шкала: например, оценка «4» = 110 % выплаты, "
            "оценка «5» = 120 %. Обе шкалы открыты в Положении.",
         7.15, 3.05, 4.9, 0.95, size=13.5, color=INK, font="Calibri",
         line_spacing=1.08)

    # Нижняя полоса: на кого распространяется + оценки
    rect(s, 0.95, 4.35, 11.4, 1.95, TINT2, line=LINE, rounded=True, radius=0.04)
    text(s, "На кого распространяется", 1.35, 4.6, 5.0, 0.4, size=13,
         bold=True, color=TEAL_D, font="Calibri")
    bullets(s, [
        "ИТ-работники бэндов M, LM и S",
        "Единые правила для БЕ Технологии, Цифровых Платформ и БЕ Роботикс",
        "Итоговая оценка эффективности — по шкале от 1 до 5",
    ], 1.35, 5.05, 6.0, 1.2, size=13.5, glyph="—", glyph_color=TEAL)

    # стат-выноска справа
    badge(s, 8.7, 4.85, 1.15, DARK, "1–5", GOLD, glyph_size=26, font="Cambria")
    text(s, "оценка\nэффективности", 10.05, 4.85, 2.2, 1.2, size=14, bold=True,
         color=INK, font="Calibri", anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
    return s


def s04_evaluation(prs):
    s = blank_slide(prs)
    set_bg(s, WHITE)
    slide_title(s, "Как формируется оценка", "Прозрачный путь от ваших КПЭ к премии")

    steps = [
        ("Ваши КПЭ и достижения", "То, что вы запланировали и сделали за полугодие"),
        ("Селфревью", "Вы самостоятельно оцениваете свои результаты"),
        ("Оценка руководителя", "Руководитель формирует свою оценку"),
        ("Калибровочная сессия", "Объективная оценка — выверяют несколько менеджеров"),
        ("Итоговая оценка 1–5", "Финальный результат за полугодие"),
        ("Премия", "Рассчитывается на основе итоговой оценки"),
    ]
    top0 = 1.95
    row_h = 0.83
    x = 1.05
    for i, (title, sub) in enumerate(steps):
        y = top0 + i * row_h
        is_last = (i == len(steps) - 1)
        b_fill = GOLD if is_last else TEAL
        badge(s, x, y, 0.56, b_fill, str(i + 1), WHITE, glyph_size=18,
              font="Cambria")
        # соединительная линия между кружками
        if not is_last:
            conn = rect(s, x + 0.27 - 0.015, y + 0.56, 0.03, row_h - 0.56,
                        LINE)
        text(s, title, x + 0.85, y - 0.04, 6.5, 0.4, size=16, bold=True,
             color=INK if not is_last else TEAL_D, font="Calibri",
             anchor=MSO_ANCHOR.MIDDLE)
        text(s, sub, x + 0.85, y + 0.33, 8.0, 0.4, size=12.5, color=MUTED,
             font="Calibri")
    return s


def s05_kf_intro(prs):
    s = blank_slide(prs)
    set_bg(s, WHITE)
    slide_title(s, "Корректировки на 2 полугодие", "Что изменилось и как это влияет на премию")

    # Карточка А: руководители бэнда M
    rect(s, 0.95, 1.9, 5.55, 2.35, TINT, rounded=True, radius=0.045)
    badge(s, 1.3, 2.25, 0.55, TEAL_D, "M", WHITE, glyph_size=18, font="Cambria")
    text(s, "Руководители бэнда M + ВП", 2.05, 2.27, 4.3, 0.5, size=16,
         bold=True, color=INK, font="Calibri", anchor=MSO_ANCHOR.MIDDLE)
    text(s, "Во 2-м полугодии в коэффициенте фонда учитывается выполнение "
            "корпоративных показателей.",
         1.3, 2.98, 4.95, 0.8, size=13, color=INK, font="Calibri",
         line_spacing=1.08)
    text(s, "25 % полугодовой премии — в 2 раза меньше, чем в прошлой системе",
         1.3, 3.75, 4.95, 0.5, size=12.5, bold=True, color=TEAL_D,
         font="Calibri", line_spacing=1.05)

    # Карточка Б: коэффициент фонда
    rect(s, 6.8, 1.9, 5.55, 2.35, GOLDT, rounded=True, radius=0.045)
    badge(s, 7.15, 2.25, 0.55, GOLD, "Кф", WHITE, glyph_size=15, font="Cambria")
    text(s, "Коэффициент фонда (Кф)", 7.9, 2.27, 4.3, 0.5, size=16, bold=True,
         color=INK, font="Calibri", anchor=MSO_ANCHOR.MIDDLE)
    text(s, "Мультипликатор, который влияет на всю премию и связывает её "
            "с реальными результатами бизнеса.",
         7.15, 2.98, 4.95, 1.1, size=13, color=INK, font="Calibri",
         line_spacing=1.08)

    # Нижний блок: когда применяется
    text(s, "Когда применяется коэффициент фонда", 0.95, 4.55, 11, 0.45,
         size=15, bold=True, color=DARK, font="Cambria")
    cards = [
        ("Бюджет фонда", "Превышение бюджетного фонда подразделения"),
        ("Финансы компании", "Текущее финансовое состояние Компании"),
        ("Корпоративные КПЭ", "Степень выполнения корпоративных КПЭ"),
    ]
    cw, gap = 3.6, 0.3
    x0 = 0.95
    for i, (t_, d_) in enumerate(cards):
        x = x0 + i * (cw + gap)
        rect(s, x, 5.05, cw, 1.4, TINT2, line=LINE, rounded=True, radius=0.06)
        oval(s, x + 0.3, 5.35, 0.22, GOLD)
        text(s, t_, x + 0.65, 5.27, cw - 0.8, 0.4, size=14, bold=True,
             color=INK, font="Calibri")
        text(s, d_, x + 0.3, 5.75, cw - 0.55, 0.6, size=12.5, color=MUTED,
             font="Calibri", line_spacing=1.05)
    return s


def s06_kf_scenarios(prs):
    s = blank_slide(prs)
    set_bg(s, WHITE)
    slide_title(s, "Коэффициент фонда: как это работает",
                "Итог зависит от бюджета и корпоративных показателей")

    scen = [
        ("Повышающий", TEAL,
         "Бюджет в норме (Кф = 1)",
         "Корп. показатели перевыполнены",
         "Итог: Кф повышающий"),
        ("Понижающий", GOLD,
         "Бюджет превышен (Кф < 1)",
         "Корп. показатели на 100 %",
         "Итог: Кф понижающий"),
        ("Понижающий", GOLD,
         "Бюджет в норме (Кф = 1)",
         "Корп. показатели частично",
         "Итог: Кф понижающий"),
    ]
    cw, gap = 3.66, 0.3
    x0 = 0.95
    top = 1.95
    ch = 2.95
    for i, (label, accent, line1, line2, res) in enumerate(scen):
        x = x0 + i * (cw + gap)
        rect(s, x, top, cw, ch, TINT2, line=LINE, rounded=True, radius=0.05)
        badge(s, x + 0.3, top + 0.32, 0.5, accent,
              "↑" if label == "Повышающий" else "↓", WHITE, glyph_size=20)
        text(s, label, x + 0.95, top + 0.32, cw - 1.1, 0.5, size=16, bold=True,
             color=INK, font="Calibri", anchor=MSO_ANCHOR.MIDDLE)
        text(s, "Бюджет фонда", x + 0.32, top + 1.05, cw - 0.6, 0.3, size=11,
             bold=True, color=MUTED, font="Calibri")
        text(s, line1, x + 0.32, top + 1.32, cw - 0.6, 0.4, size=13, color=INK,
             font="Calibri")
        text(s, "Корп. показатели", x + 0.32, top + 1.75, cw - 0.6, 0.3,
             size=11, bold=True, color=MUTED, font="Calibri")
        text(s, line2, x + 0.32, top + 2.02, cw - 0.6, 0.4, size=13, color=INK,
             font="Calibri")
        # итог-плашка внизу карточки
        rect(s, x + 0.32, top + 2.5, cw - 0.64, 0.34, accent, rounded=True,
             radius=0.4)
        text(s, res, x + 0.32, top + 2.5, cw - 0.64, 0.34, size=11.5,
             bold=True, color=WHITE, font="Calibri", align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)

    # Нижняя выноска-цитата
    rect(s, 0.95, 5.2, 11.4, 1.1, DARK, rounded=True, radius=0.06)
    badge(s, 1.3, 5.5, 0.5, GOLD, "!", WHITE, glyph_size=20)
    text(s, "Мы ввели Кф, чтобы результаты компании влияли справедливо на всех: "
            "если бизнес растёт — вырастет и ваша премия. Если есть сложности — "
            "мы вместе сокращаем.",
         2.05, 5.35, 10.0, 0.85, size=13.5, color=WHITE, font="Calibri",
         italic=True, line_spacing=1.1, anchor=MSO_ANCHOR.MIDDLE)
    return s


def s07_example_calc(prs):
    s = blank_slide(prs)
    set_bg(s, WHITE)
    slide_title(s, "Практический пример: считаем вместе",
                "Как из оклада получается итоговая премия")

    # Левая колонка — исходные данные
    rect(s, 0.95, 1.95, 4.3, 4.35, DARK, rounded=True, radius=0.04)
    text(s, "Исходные данные", 1.3, 2.25, 3.7, 0.45, size=16, bold=True,
         color=GOLD, font="Cambria")
    data = [
        ("Оклад", "150 000 ₽ / мес."),
        ("Целевой % премии", "30 %"),
        ("Оценка эффективности", "4  →  %ВП = 110 %"),
        ("Коэффициент фонда", "1.0"),
    ]
    yy = 2.95
    for k, v in data:
        text(s, k, 1.3, yy, 3.7, 0.3, size=12, color=TINT, font="Calibri")
        text(s, v, 1.3, yy + 0.28, 3.7, 0.4, size=16, bold=True, color=WHITE,
             font="Calibri")
        yy += 0.82

    # Правая колонка — пошаговый расчёт
    steps = [
        ("Шаг 1 · Окладный фонд", "150 000 × 6 мес.", "900 000 ₽"),
        ("Шаг 2 · Целевая сумма премии", "900 000 × 30 %", "270 000 ₽"),
        ("Шаг 3 · Применяем %ВП (оценка 4 → 110 %)", "270 000 × 110 %", "297 000 ₽"),
        ("Шаг 4 · Применяем Кф", "297 000 × 1.0", "297 000 ₽"),
    ]
    x = 5.55
    yy = 1.95
    sh = 1.0
    for i, (title, formula, result) in enumerate(steps):
        rect(s, x, yy, 6.8, 0.85, TINT2 if i % 2 == 0 else TINT, line=LINE,
             rounded=True, radius=0.06)
        badge(s, x + 0.25, yy + 0.22, 0.42, TEAL, str(i + 1), WHITE,
              glyph_size=15, font="Cambria")
        text(s, title, x + 0.85, yy + 0.13, 4.0, 0.35, size=13, bold=True,
             color=INK, font="Calibri")
        text(s, formula, x + 0.85, yy + 0.46, 4.0, 0.32, size=12, color=MUTED,
             font="Calibri")
        text(s, result, x + 4.7, yy, 1.95, 0.85, size=18, bold=True,
             color=TEAL_D, font="Calibri", align=PP_ALIGN.RIGHT,
             anchor=MSO_ANCHOR.MIDDLE)
        yy += sh
    return s


def s08_result(prs):
    s = blank_slide(prs)
    set_bg(s, DARK)
    # мотив
    oval(s, 10.6, 4.6, 4.0, TEAL)
    oval(s, 12.2, 5.9, 2.0, GOLD)

    text(s, "Ваша полугодовая премия", 0.95, 1.4, 9, 0.6, size=20, color=TINT,
         font="Calibri")
    text(s, "297 000 ₽", 0.9, 1.95, 9.5, 1.7, size=96, bold=True, color=GOLD,
         font="Cambria")
    text(s, "вместо 270 000 ₽ целевых  —  +27 000 ₽ за оценку «4»",
         0.97, 3.85, 9.5, 0.5, size=18, color=WHITE, font="Calibri")

    # карточка-вывод
    rect(s, 0.95, 4.7, 8.6, 1.6, TEAL_D, rounded=True, radius=0.05)
    badge(s, 1.3, 5.05, 0.55, GOLD, "✓", WHITE, glyph_size=22)
    text(s, "Это ваши деньги за ваши результаты", 2.1, 5.0, 7.2, 0.45, size=16,
         bold=True, color=WHITE, font="Calibri")
    text(s, "Оценка «4» даёт +10 % сверх целевой суммы. Оценка «5» дала бы +20 %.",
         2.1, 5.5, 7.2, 0.6, size=13.5, color=TINT, font="Calibri",
         line_spacing=1.1)
    return s


def s09_guarantees(prs):
    s = blank_slide(prs)
    set_bg(s, WHITE)
    slide_title(s, "Гарантии для вас", "Мы позаботились о вашей защите")

    cards = [
        ("%", GOLD, "Минимальная премия — 10 %",
         "Даже при оценке «1» вы получите не менее 10 % от целевой суммы. "
         "Закреплено прямо в Шкалах."),
        ("=", TEAL, "Объективная оценка",
         "Калибровочная сессия исключает субъективизм одного руководителя — "
         "оценку выверяют несколько менеджеров."),
        ("□", TEAL, "Открытые шкалы",
         "Формула и шкалы №1 и №2 — прямо в Положении. Никаких «чёрных ящиков». "
         "Доп. ознакомление в Сфере."),
        ("♥", GOLD, "Защита декретных работников",
         "В отпуске по БиР или уходу за ребёнком право на премию сохраняется."),
    ]
    cw, ch, gx, gy = 5.55, 1.55, 0.3, 0.28
    x0, y0 = 0.95, 1.85
    for i, (glyph, accent, title, desc) in enumerate(cards):
        col = i % 2
        row = i // 2
        x = x0 + col * (cw + gx)
        y = y0 + row * (ch + gy)
        rect(s, x, y, cw, ch, TINT2, line=LINE, rounded=True, radius=0.05)
        badge(s, x + 0.32, y + 0.32, 0.55, accent, glyph, WHITE, glyph_size=20)
        text(s, title, x + 1.1, y + 0.3, cw - 1.3, 0.55, size=15.5, bold=True,
             color=INK, font="Calibri", anchor=MSO_ANCHOR.MIDDLE)
        text(s, desc, x + 0.35, y + 0.92, cw - 0.65, 0.55, size=12.5,
             color=MUTED, font="Calibri", line_spacing=1.05)

    # цитата
    text(s, "«Система создана для того, чтобы вы были вознаграждены за реальные "
            "результаты, а не за „видимость работы\". Открытые шкалы и калибровка "
            "— это инструменты честности.»",
         0.95, 6.05, 11.4, 0.8, size=13, italic=True, color=TEAL_D,
         font="Cambria", align=PP_ALIGN.CENTER, line_spacing=1.08)
    return s


def s10_timeline(prs):
    s = blank_slide(prs)
    set_bg(s, WHITE)
    slide_title(s, "Сроки и условия выплаты", "Когда и при каких условиях")

    # Таймлайн — два периода
    rect(s, 0.95, 1.95, 5.55, 1.5, TINT, rounded=True, radius=0.05)
    badge(s, 1.3, 2.3, 0.55, TEAL, "H1", WHITE, glyph_size=15, font="Cambria")
    text(s, "1-е полугодие (янв–июн)", 2.05, 2.28, 4.2, 0.45, size=15,
         bold=True, color=INK, font="Calibri", anchor=MSO_ANCHOR.MIDDLE)
    text(s, "Выплата премии за H1 — сентябрь текущего года", 1.3, 2.95, 4.95,
         0.4, size=13, color=MUTED, font="Calibri")

    rect(s, 6.8, 1.95, 5.55, 1.5, TINT, rounded=True, radius=0.05)
    badge(s, 7.15, 2.3, 0.55, GOLD, "H2", WHITE, glyph_size=15, font="Cambria")
    text(s, "2-е полугодие (июл–дек)", 7.9, 2.28, 4.2, 0.45, size=15,
         bold=True, color=INK, font="Calibri", anchor=MSO_ANCHOR.MIDDLE)
    text(s, "Выплата премии за H2 — март следующего года", 7.15, 2.95, 4.95,
         0.4, size=13, color=MUTED, font="Calibri")

    # Условия получения
    text(s, "Условия получения премии", 0.95, 3.85, 11, 0.45, size=15,
         bold=True, color=DARK, font="Cambria")
    conds = [
        ("Минимальный стаж", "30 рабочих дней в полугодии"),
        ("Карта КПЭ", "Установлена и согласована"),
        ("Оценка эффективности", "Сформирована на калибровочной сессии"),
        ("Трудовые отношения", "Действуют на дату формирования приказа"),
    ]
    yy = 4.4
    rh = 0.6
    for i, (k, v) in enumerate(conds):
        rect(s, 0.95, yy, 11.4, 0.5, TINT2 if i % 2 == 0 else WHITE,
             line=LINE if i % 2 else None, rounded=True, radius=0.08)
        oval(s, 1.25, yy + 0.14, 0.22, TEAL)
        text(s, k, 1.7, yy, 4.0, 0.5, size=13.5, bold=True, color=INK,
             font="Calibri", anchor=MSO_ANCHOR.MIDDLE)
        text(s, v, 5.8, yy, 6.4, 0.5, size=13.5, color=MUTED, font="Calibri",
             anchor=MSO_ANCHOR.MIDDLE)
        yy += rh
    return s


def s11_faq(prs):
    s = blank_slide(prs)
    set_bg(s, WHITE)
    slide_title(s, "Частые вопросы", "Отвечаем на то, что спрашивают чаще всего")

    qa = [
        ("Я в отпуске по уходу за ребёнком — мне ничего не заплатят?",
         "Право на премию сохраняется, если выполнены условия Положения "
         "(КПЭ установлены, оценка проведена, отработано 30 рабочих дней). "
         "Вы получите премию за фактически отработанное время."),
        ("Что такое калибровочная сессия и почему я в ней не участвую?",
         "Это встреча руководителей, где они обсуждают и выравнивают оценки "
         "сотрудников на сопоставимых ролях. Так у двух одинаковых специалистов "
         "из разных команд оценки не отличаются «из-за характера руководителя». "
         "Ваш руководитель потом даёт вам обратную связь."),
    ]
    y0 = 1.95
    ch = 2.05
    gap = 0.3
    for i, (q, a) in enumerate(qa):
        y = y0 + i * (ch + gap)
        rect(s, 0.95, y, 11.4, ch, TINT2, line=LINE, rounded=True, radius=0.04)
        badge(s, 1.3, y + 0.32, 0.55, GOLD, "?", WHITE, glyph_size=22)
        text(s, q, 2.1, y + 0.3, 9.9, 0.55, size=16, bold=True, color=INK,
             font="Calibri", anchor=MSO_ANCHOR.MIDDLE)
        text(s, a, 2.1, y + 0.95, 9.9, 0.95, size=13.5, color=MUTED,
             font="Calibri", line_spacing=1.1)
    return s


def s12_next_steps(prs):
    s = blank_slide(prs)
    set_bg(s, WHITE)
    slide_title(s, "Ваши следующие шаги", "Что сделать прямо сейчас")

    steps = [
        ("Ознакомьтесь с Положением",
         "Пришлём на ознакомление в Сферу. Обратите внимание на Шкалы (Приложение №1)."),
        ("Обсудите с руководителем цели на 2-е полугодие",
         "Внесите их в систему, проконтролируйте согласование, следите за дедлайнами КПЭ."),
        ("Собирайте факты для оценки",
         "Фиксируйте достижения, метрики и результаты — пригодится для самооценки и калибровки."),
        ("Задайте вопросы",
         "Управление по компенсациям и льготам · ваш HR-партнёр."),
    ]
    y0 = 1.9
    rh = 0.92
    for i, (title, sub) in enumerate(steps):
        y = y0 + i * rh
        badge(s, 1.05, y, 0.6, TEAL, str(i + 1), WHITE, glyph_size=18,
              font="Cambria")
        text(s, title, 1.9, y - 0.02, 6.3, 0.4, size=16, bold=True, color=INK,
             font="Calibri", anchor=MSO_ANCHOR.MIDDLE)
        text(s, sub, 1.9, y + 0.37, 6.8, 0.45, size=12.5, color=MUTED,
             font="Calibri", line_spacing=1.05)

    # Контактная карточка справа
    rect(s, 9.1, 1.9, 3.25, 3.7, DARK, rounded=True, radius=0.05)
    oval(s, 11.3, 4.6, 1.8, TEAL)
    badge(s, 9.5, 2.3, 0.62, GOLD, "@", WHITE, glyph_size=22)
    text(s, "Мы на связи!", 9.5, 3.1, 2.6, 0.5, size=20, bold=True,
         color=WHITE, font="Cambria")
    text(s, "Управление по компенсациям и льготам", 9.5, 3.75, 2.55, 0.7,
         size=12.5, color=TINT, font="Calibri", line_spacing=1.1)
    text(s, "Ваш HR-партнёр", 9.5, 4.5, 2.55, 0.4, size=12.5, color=GOLD,
         bold=True, font="Calibri")
    return s


# ----------------------------------------------------------------------------
# Сборка презентации
# ----------------------------------------------------------------------------
def build(output_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)   # формат 16:9
    prs.slide_height = Inches(7.5)

    s01_title(prs)
    s02_contents(prs)
    s03_how_it_works(prs)
    s04_evaluation(prs)
    s05_kf_intro(prs)
    s06_kf_scenarios(prs)
    s07_example_calc(prs)
    s08_result(prs)
    s09_guarantees(prs)
    s10_timeline(prs)
    s11_faq(prs)
    s12_next_steps(prs)

    prs.save(output_path)
    return output_path


def main():
    # Папка из аргумента командной строки или значение по умолчанию
    out_dir = sys.argv[1] if len(sys.argv) > 1 else DEFAULT_OUTPUT_DIR
    out_dir = Path(out_dir).expanduser()
    out_dir.mkdir(parents=True, exist_ok=True)

    out_path = out_dir / OUTPUT_FILENAME
    build(str(out_path))
    print("Готово! Презентация сохранена:")
    print("   " + str(out_path.resolve()))


if __name__ == "__main__":
    main()
